from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

ACCENT = RGBColor(0, 102, 204)
DARK = RGBColor(0x1a, 0x1a, 0x2e)
LIGHT_BG = RGBColor(0xf5, 0xf7, 0xfa)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x52, 0xc4, 0x1a)
RED = RGBColor(0xff, 0x4d, 0x4f)
ORANGE = RGBColor(0xfa, 0x8c, 0x16)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title, subtitle=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.2)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xcc, 0xdd, 0xff)
        p2.alignment = PP_ALIGN.LEFT

def add_body_text(slide, text, left, top, width, height, size=14, bold=False, color=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or DARK
    return tf

def add_bullet_list(slide, items, left, top, width, height, size=14):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)
    return tf

def add_page_number(slide, num, total):
    txBox = slide.shapes.add_textbox(Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f'{num}/{total}'
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT

TOTAL = 18

# ==================== SLIDE 1: COVER ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)
add_body_text(slide, 'Tokenhub', Inches(1), Inches(1.5), Inches(11), Inches(1.2), size=52, bold=True, color=WHITE)
add_body_text(slide, '企业AI治理智能平台', Inches(1), Inches(2.7), Inches(11), Inches(0.8), size=32, color=ACCENT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.6), Inches(3), Inches(0.05))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()
add_body_text(slide, '商业计划书  |  天使轮  |  ¥500万', Inches(1), Inches(4.0), Inches(11), Inches(0.5), size=18, color=GRAY)
add_body_text(slide, '2026年6月', Inches(1), Inches(4.8), Inches(11), Inches(0.5), size=14, color=GRAY)
add_page_number(slide, 1, TOTAL)

# ==================== SLIDE 2: ONE-LINER ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '项目概述', 'Project Overview')

add_body_text(slide, '一句话定位', Inches(0.8), Inches(1.6), Inches(11), Inches(0.4), size=20, bold=True, color=ACCENT)
add_body_text(slide, 'Tokenhub 是面向企业客户的 AI 治理与 Token 安全运营中台，帮助企业安全、合规、可控地使用大模型，同时降低 30%+ AI 调用成本。',
              Inches(0.8), Inches(2.1), Inches(11), Inches(0.8), size=16)

add_body_text(slide, '产品形态', Inches(0.8), Inches(3.2), Inches(11), Inches(0.4), size=20, bold=True, color=ACCENT)
add_bullet_list(slide, [
    '• SaaS + 私有化部署双模式，已接入 8 家模型厂商（DeepSeek/通义千问/豆包/GLM/Kimi/文心/混元等）',
    '• 提供统一API网关、智能路由、安全过滤、预算管控、合规审计、Agent市场、知识库治理等全栈能力',
], Inches(0.8), Inches(3.7), Inches(11), Inches(1), size=14)

add_body_text(slide, '当前阶段', Inches(0.8), Inches(4.9), Inches(11), Inches(0.4), size=20, bold=True, color=ACCENT)
add_bullet_list(slide, [
    '• MVP 已完成：15个后端模块、70+ API、14个前端页面，可现场演示',
    '• 测试覆盖：Go单测30+、前端40+、E2E 6+，全部通过',
], Inches(0.8), Inches(5.4), Inches(11), Inches(0.8), size=14)

add_page_number(slide, 2, TOTAL)

# ==================== SLIDE 3: MARKET SIZE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '市场机会与规模', 'Market Opportunity')

rows_data = [
    ['市场领域', '2024规模', '2028预测', 'CAGR', '来源'],
    ['全球AI治理市场', '$1.8B', '$6.2B', '28.2%', 'MarketsandMarkets 2025'],
    ['中国AI安全市场', '¥48亿', '¥186亿', '31.1%', 'IDC China 2025'],
    ['API管理/网关', '$5.4B', '$17.3B', '26.1%', 'Gartner 2025'],
    ['中国AI平台市场', '¥132亿', '¥508亿', '30.9%', '艾瑞咨询 2025'],
    ['AI Agent市场', '$5.1B', '$47.1B', '55.8%', 'MarketsandMarkets 2025'],
]

table_shape = slide.shapes.add_table(len(rows_data), len(rows_data[0]), Inches(0.5), Inches(1.6), Inches(12.3), Inches(3.5))
table = table_shape.table

for r in range(len(rows_data)):
    for c in range(len(rows_data[0])):
        cell = table.cell(r, c)
        cell.text = rows_data[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT
            else:
                p.font.color.rgb = DARK
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_BG
            p.alignment = PP_ALIGN.CENTER

add_body_text(slide, 'TAM/SAM/SOM 测算', Inches(0.8), Inches(5.4), Inches(11), Inches(0.4), size=16, bold=True, color=ACCENT)
add_bullet_list(slide, [
    '• TAM（总市场）：中国AI治理与安全市场 — ¥186亿（2028E）',
    '• SAM（可服务）：中大型企业AI治理平台 — ¥37亿',
    '• SOM（可获得）：天使轮目标 — 前3年签约100家客户，ARR ¥5,000万',
], Inches(0.8), Inches(5.9), Inches(11), Inches(1), size=13)

add_page_number(slide, 3, TOTAL)

# ==================== SLIDE 4: PAIN POINTS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '行业痛点：企业AI落地的6大困境', 'Industry Pain Points')

pain_points = [
    ('模型碎片化', '多厂商、多模型、多Key分散管理'),
    ('成本黑洞', 'Token消耗不透明，预算频繁超支'),
    ('安全风险', '数据泄露/提示注入/越权调用'),
    ('Agent失控', '自主循环消耗、工具链投毒'),
    ('合规压力', '算法备案/等保/数据分类分级'),
    ('计量纠纷', '多方对账口径不一致，结算争议'),
]

for i, (title, desc) in enumerate(pain_points):
    col = i % 3
    row = i // 3
    left = 0.8 + col * 4.2
    top = 1.8 + row * 2.6
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.8), Inches(2.2))
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_BG; shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = f'0{i+1}  {title}'
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13); p2.font.color.rgb = DARK

add_page_number(slide, 4, TOTAL)

# ==================== SLIDE 5: PRODUCT ARCH ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '产品架构：三层治理模型', 'Product Architecture')

layers = [
    ('消费侧：企业内部治理', '统一接入网关 | 凭证托管 | 场景登记 | 数据分级 | AI FinOps', GREEN),
    ('平台侧：安全运营中枢', 'Token生命周期 | 计量账本 | 安全网关 | 智能路由 | 预算熔断 | 知识库安全', ACCENT),
    ('供给侧：模型聚合与交付', '8厂商9模型 | 统一API | 智能路由 | Agent市场 | 服务目录', ORANGE),
]

for i, (title, desc, color) in enumerate(layers):
    top = 1.6 + i * 1.8
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top), Inches(11.5), Inches(1.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.5); tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = title; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = desc; p2.font.size = Pt(14); p2.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

add_page_number(slide, 5, TOTAL)

# ==================== SLIDE 6: FEATURE MATRIX ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '核心功能矩阵', 'Feature Matrix')

features = [
    ['模块', '功能亮点', '状态'],
    ['模型方舟', '8厂商9模型统一接入 + 智能路由 + 场景推荐', '✅ 已完成'],
    ['安全网关', '敏感词14类 + PII脱敏 + 注入检测 + DLP', '✅ 已完成'],
    ['预算管控', '5级预算 + 3态熔断 + Agent Guard', '✅ 已完成'],
    ['安全巡检Agent', '注入/权限/合规/供应链 4维自动扫描', '✅ 已完成'],
    ['智能体市场', '专家/Skills/Connectors/MCP工具生态', '✅ 已完成'],
    ['知识库管理', '上传即扫描，安全联动（4维度检测）', '✅ 已完成'],
    ['合规报告', '使用台账/安全审计/算法备案 三大报告', '✅ 已完成'],
    ['Token Broker', '短期动态凭证 + 即时撤销', '📋 V1.0规划'],
    ['计量账本', '不可篡改交易数据底座', '📋 V1.0规划'],
]

table_shape = slide.shapes.add_table(len(features), 3, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
table = table_shape.table
for r in range(len(features)):
    for c in range(3):
        cell = table.cell(r, c)
        cell.text = features[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            if r == 0:
                p.font.bold = True; p.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
            else:
                p.font.color.rgb = DARK
                if r % 2 == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_BG
            p.alignment = PP_ALIGN.CENTER

add_page_number(slide, 6, TOTAL)

# ==================== SLIDE 7: KEY FEATURE - SAFETY ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '产品亮点 ①：四层安全引擎', 'Safety Engine')

security_items = [
    ('第1层：敏感词检测', 'Aho-Corasick 自动机，14类敏感词（涉政/色情/暴力/赌博/违禁品\n/诈骗/暴恐/未成年人保护等），毫秒级实时匹配，支持热更新'),
    ('第2层：PII 检测与脱敏', '自动识别身份证号、手机号、邮箱、银行卡号\n支持脱敏处理（如 138****1234），防止敏感数据泄露'),
    ('第3层：提示注入检测', '检测 DAN 越狱攻击、系统角色逃逸、Token 泄露\n代码注入、数据窃取等高级攻击模式'),
    ('第4层：DLP 安全网关', '6条网关规则（敏感词/PII/注入/越权/预算/限流）\n统一身份认证 + SSO（LDAP/钉钉/企业微信/OAuth）+ IP白名单'),
]

for i, (title, desc) in enumerate(security_items):
    top = 1.5 + i * 1.45
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top), Inches(11.5), Inches(1.2))
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_BG; shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = desc; p2.font.size = Pt(12); p2.font.color.rgb = DARK

add_page_number(slide, 7, TOTAL)

# ==================== SLIDE 8: KEY FEATURE - BUDGET ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '产品亮点 ②：五级预算管控 + Agent Guard', 'Budget Control')

add_bullet_list(slide, [
    '五级预算体系：公司 → 部门 → 项目 → 用户 → API Key，每一级独立预算上限',
    '',
    '三级熔断器（基于 Redis Lua 原子操作）：',
    '  • 使用率 ≥ 80% → throttled（限速至10%）',
    '  • 使用率 ≥ 90% → throttled（限速至5%）',
    '  • 使用率 ≥ 100% → blocked（24小时硬阻断）',
    '',
    'Agent Guard 三项异常检测规则：',
    '  • cost_spike：5分钟内花费 > ¥500 → 自动熔断',
    '  • loop_detection：1分钟内调用 > 50次 → 触发限速',
    '  • same_prompt_loop：相同提示词重复 > 10次 → 阻断15分钟',
], Inches(0.8), Inches(1.5), Inches(11), Inches(5.5), size=15)

add_page_number(slide, 8, TOTAL)

# ==================== SLIDE 9: KEY FEATURE - AGENT ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '产品亮点 ③：安全巡检 Agent（行业首创）', 'Security Agent')

add_bullet_list(slide, [
    '自动对 Agent、Skill、MCP 工具进行四维度安全扫描：',
    '',
    '  🔍 注入扫描器：检测提示注入、命令注入、硬编码凭据（8条规则）',
    '  🔐 权限审计器：API Key 暴露、过度权限、预算过高等',
    '  📋 合规检查器：敏感词命中、PII 泄露等',
    '  🔗 供应链扫描器：第三方连接器风险、MCP 配置安全',
    '',
    '触发机制：',
    '  • 资产变更自动扫描（10秒防抖）',
    '  • 24小时定时全量扫描',
    '  • 手动一键全量扫描',
    '',
    '扫描结果自动同步到知识库条目状态（safe / risky / blocked）',
], Inches(0.8), Inches(1.5), Inches(11), Inches(5.5), size=15)

add_page_number(slide, 9, TOTAL)

# ==================== SLIDE 10: BUSINESS MODEL ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '商业模式', 'Business Model')

pricing = [
    ['版本', '价格(年)', '核心功能', '目标客户'],
    ['社区版', '免费', '10个Key + 基础安全 + 3模型', '开发者/初创'],
    ['专业版', '¥9.8万', '无限Key + 高级安全 + FinOps', '中小企业'],
    ['企业版', '¥29.8万', '私有部署 + SSO + 定制合规', '政企大客户'],
    ['旗舰版', '¥68万+', '专属Agent + 交易平台', '运营商/超大企业'],
]

table_shape = slide.shapes.add_table(len(pricing), 4, Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.5))
table = table_shape.table
for r in range(len(pricing)):
    for c in range(4):
        cell = table.cell(r, c); cell.text = pricing[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            if r == 0:
                p.font.bold = True; p.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
            else:
                p.font.color.rgb = DARK
            p.alignment = PP_ALIGN.CENTER

add_body_text(slide, '附加收入', Inches(0.8), Inches(4.5), Inches(11), Inches(0.4), size=16, bold=True, color=ACCENT)
add_bullet_list(slide, [
    '• Agent市场佣金（15-20%）• 安全评估服务（¥5-15万/次）• 定制Agent开发 • API开放平台',
], Inches(0.8), Inches(5.0), Inches(11), Inches(0.5), size=13)

income = [
    ['年份', '客户数', 'ARR(万元)', '毛利', '净利'],
    ['Y1', '30', '350', '65%', '-30%'],
    ['Y2', '80', '1,500', '72%', '8%'],
    ['Y3', '200', '5,000', '78%', '22%'],
]
table_shape2 = slide.shapes.add_table(len(income), 5, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.5))
table2 = table_shape2.table
for r in range(len(income)):
    for c in range(5):
        cell = table2.cell(r, c); cell.text = income[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            if r == 0:
                p.font.bold = True; p.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
            p.alignment = PP_ALIGN.CENTER

add_page_number(slide, 10, TOTAL)

# ==================== SLIDE 11: COMPETITION OVERVIEW ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '竞争分析：与竞品对比', 'Competitive Analysis')

comps = [
    ['维度', 'Tokenhub', '硅基流动', '天翼云息壤', '阿里云百炼'],
    ['定位', '治理+安全中台', 'API聚合+加速', '算力调度+运营', '全栈AI中台'],
    ['模型', '8厂商9模型', '30+模型', '自研为主', '自研为主'],
    ['安全', '4层引擎 ★★★★', '基础 ★', '链路安全 ★★', '内容安全 ★★'],
    ['预算', '5级熔断 ★★★★', '基础 ★', '配额 ★★', '配额 ★★'],
    ['Agent治理', '市场+巡检 ★★★★', '无', '无', '开发平台 ★★'],
    ['知识库安全', '上传即扫描 ★★★', '无', '无', '无'],
    ['合规', '三大报告 ★★★', '基础统计 ★', '报表 ★★', '基础审计 ★'],
]

table_shape = slide.shapes.add_table(len(comps), 5, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
table = table_shape.table
for r in range(len(comps)):
    for c in range(5):
        cell = table.cell(r, c); cell.text = comps[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            if r == 0:
                p.font.bold = True; p.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
            elif c == 1:  # highlight Tokenhub column
                p.font.bold = True
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xe6, 0xf0, 0xff)
            p.alignment = PP_ALIGN.CENTER

add_page_number(slide, 11, TOTAL)

# ==================== SLIDE 12: COMPETITION ADVANTAGE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '差异化优势', 'Competitive Advantage')

advantages = [
    '① 唯一全链路覆盖：同时覆盖"供给侧→平台侧→消费侧→智能体侧"，单一竞品仅1-2环节',
    '② 安全深度领先：4层引擎 vs 竞品1-2层，14类敏感词毫秒级检测',
    '③ Agent安全巡检：行业首创对Agent/Skill/MCP四维度自动扫描',
    '④ 厂商中立：不绑定模型厂商，客户自由选择最优组合',
    '⑤ 知识库安全联动：RAG场景安全盲区解决方案',
    '⑥ 产品已完成MVP：可演示，70+ API vs 竞品大多还处于概念阶段',
]

add_bullet_list(slide, advantages, Inches(0.8), Inches(1.6), Inches(11), Inches(2.5), size=16)

add_body_text(slide, '竞争壁垒', Inches(0.8), Inches(4.5), Inches(11), Inches(0.4), size=18, bold=True, color=ACCENT)
add_bullet_list(slide, [
    '• 技术壁垒：多模型适配器 + Aho-Corasick安全引擎 + Agent巡检系统',
    '• 数据壁垒：企业AI调用数据沉淀 → 安全规则持续优化 → 厂商切换成本高',
    '• 生态壁垒：Agent市场 → 开发者生态 → 网络效应',
    '• 先发优势：目前市场无同类全链路治理产品',
], Inches(0.8), Inches(5.0), Inches(11), Inches(2), size=14)

add_page_number(slide, 12, TOTAL)

# ==================== SLIDE 13: ROADMAP ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '发展路线图', 'Development Roadmap')

timeline = [
    ('2026 Q2\n当前', 'MVP验证期', '产品可演示\n种子用户试用\n10家试用客户'),
    ('2026 Q3-Q4\nV1.0', '商用版发布', 'Token Broker\n计量账本\n30家付费客户'),
    ('2027 Q1-Q2', '生态扩展', '交易平台\nAPI开放\n80家客户'),
    ('2027 Q3-Q4', '行业深耕', '金融合规版\n政务定制版\n200家客户'),
    ('2028+', '规模化增长', '生态运营商\n国际市场\nIPO准备'),
]

for i, (time, phase, milestones) in enumerate(timeline):
    left = 0.5 + i * 2.5
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.6), Inches(2.3), Inches(3.0))
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_BG; shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]; p.text = time; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = ACCENT; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = phase; p2.font.size = Pt(14); p2.font.bold = True; p2.font.color.rgb = DARK; p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph(); p3.text = milestones; p3.font.size = Pt(11); p3.font.color.rgb = DARK; p3.alignment = PP_ALIGN.CENTER

    if i < len(timeline) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left + 2.3), Inches(2.8), Inches(0.2), Inches(0.3))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT; arrow.line.fill.background()

add_body_text(slide, 'KPI 目标', Inches(0.8), Inches(5.0), Inches(11), Inches(0.4), size=16, bold=True, color=ACCENT)
kpi_data = [
    ['', '6个月', '12个月', '24个月', '36个月'],
    ['付费客户', '10', '30', '80', '200'],
    ['ARR', '¥100万', '¥350万', '¥1,500万', '¥5,000万'],
    ['团队', '8人', '15人', '30人', '50人'],
]
table_shape = slide.shapes.add_table(len(kpi_data), 5, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.8))
table = table_shape.table
for r in range(len(kpi_data)):
    for c in range(5):
        cell = table.cell(r, c); cell.text = kpi_data[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            if r == 0:
                p.font.bold = True; p.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
            p.alignment = PP_ALIGN.CENTER

add_page_number(slide, 13, TOTAL)

# ==================== SLIDE 14: TECH ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '技术亮点', 'Technical Highlights')

tech_stack = [
    ['层次', '技术选型', '核心能力'],
    ['后端', 'Go 1.22 + Gin v1.10', '高并发网关，支持万级QPS'],
    ['前端', 'React 19 + Next.js 15 + Ant Design 5', '现代化管理后台'],
    ['安全引擎', 'Aho-Corasick 自动机', '14类敏感词毫秒级匹配'],
    ['预算控制', 'Redis Lua 原子操作', '三级熔断器毫秒级判断'],
    ['Agent巡检', '事件驱动 + 4并行技能', '配置变更10秒自动扫描'],
    ['数据存储', 'PostgreSQL + ClickHouse + Redis', 'OLTP + OLAP + 缓存三层'],
    ['测试', 'Go test + Vitest + Playwright', '30+40+6 测试全通过'],
]

table_shape = slide.shapes.add_table(len(tech_stack), 3, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
table = table_shape.table
for r in range(len(tech_stack)):
    for c in range(3):
        cell = table.cell(r, c); cell.text = tech_stack[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            if r == 0:
                p.font.bold = True; p.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
            p.alignment = PP_ALIGN.CENTER

add_page_number(slide, 14, TOTAL)

# ==================== SLIDE 15: TEAM ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '团队介绍', 'Team')

add_body_text(slide, '注：以下为模板框架，请自行填充团队成员信息', Inches(0.8), Inches(1.4), Inches(11), Inches(0.4), size=12, color=GRAY)

team_roles = [
    ('创始人 / CEO', '[姓名]', '教育背景：[请填写]\n核心经验：[请填写工作经历和成就]\n愿景：[请填写创始动机与愿景]'),
    ('联合创始人 / CTO', '[姓名]', '教育背景：[请填写]\n技术能力：[请填写全栈/Go/React/安全经验]\n过往项目：[请填写主导的技术项目]'),
    ('产品负责人', '[姓名]', '教育背景：[请填写]\n行业经验：[请填写AI/企业服务产品经验]\nTaaS标准参与：[请填写行业标准贡献]'),
    ('技术顾问', '[姓名]', '教育背景：[请填写]\n行业影响力：[请填写学术/行业背景]'),
]

for i, (role, name, desc) in enumerate(team_roles):
    top = 2.0 + i * 1.3
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top), Inches(11.5), Inches(1.1))
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_BG; shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = f'{role}：{name}'; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = desc; p2.font.size = Pt(11); p2.font.color.rgb = DARK

add_page_number(slide, 15, TOTAL)

# ==================== SLIDE 16: FUNDING ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '融资需求', 'Funding Requirements')

funding_info = [
    ['融资轮次', '天使轮', '融资金额', '¥500万'],
    ['出让股权', '10-15%（待协商）', '投前估值', '¥3,300-5,000万'],
    ['资金到位', '2026年Q3', 'Runway', '18个月'],
]

table_shape = slide.shapes.add_table(len(funding_info), 4, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.5))
table = table_shape.table
for r in range(len(funding_info)):
    for c in range(4):
        cell = table.cell(r, c); cell.text = funding_info[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13); p.alignment = PP_ALIGN.CENTER
            if c == 0 or c == 2:
                p.font.bold = True
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_BG

add_body_text(slide, '资金用途', Inches(0.8), Inches(3.4), Inches(11), Inches(0.4), size=18, bold=True, color=ACCENT)
usage = [
    ['用途', '金额', '占比', '说明'],
    ['产品研发', '¥200万', '40%', 'Token Broker / 计量账本 / RAG安全\n+3后端 +2前端'],
    ['市场拓展', '¥150万', '30%', '首批10家标杆客户\n行业会议 + 销售团队（+2BD）'],
    ['合规认证', '¥75万', '15%', '等保三级 / ISO27001 / 算法备案'],
    ['运营储备', '¥75万', '15%', '云服务(18月) / 办公 / 现金流安全垫'],
]
table_shape2 = slide.shapes.add_table(len(usage), 4, Inches(0.5), Inches(3.9), Inches(12.3), Inches(2.8))
table2 = table_shape2.table
for r in range(len(usage)):
    for c in range(4):
        cell = table2.cell(r, c); cell.text = usage[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            if r == 0:
                p.font.bold = True; p.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
            p.alignment = PP_ALIGN.CENTER

add_page_number(slide, 16, TOTAL)

# ==================== SLIDE 17: INVESTMENT RETURN ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, '投资回报与退出路径', 'Investment Return & Exit')

add_body_text(slide, '回报预期', Inches(0.8), Inches(1.6), Inches(11), Inches(0.4), size=20, bold=True, color=ACCENT)
add_bullet_list(slide, [
    '• 3年后 ARR 目标：¥5,000万，对应估值 ¥3.75-5亿（7.5-10x P/S）',
    '• 天使轮投资人预期回报：3年内 10-15x',
    '• 当前可比公司估值：AI安全类 SaaS 公司 P/S 倍数 8-15x',
], Inches(0.8), Inches(2.1), Inches(11), Inches(1.5), size=15)

add_body_text(slide, '退出路径', Inches(0.8), Inches(4.0), Inches(11), Inches(0.4), size=20, bold=True, color=ACCENT)

exit_items = [
    ('A轮融资', '2027-2028', 'ARR ¥1,500万+ → 估值 ¥1.5-2亿，引入VC'),
    ('战略并购', '2028-2029', '大型云厂商/安全厂商收购，退出倍数 5-10x'),
    ('IPO', '2029+', '科创板/北交所，ARR ¥1.5亿+，估值 ¥15-20亿'),
]

for i, (title, time, desc) in enumerate(exit_items):
    top = 4.5 + i * 0.9
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top), Inches(11.5), Inches(0.8))
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_BG; shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = f'{title}'; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = desc; p2.font.size = Pt(12); p2.font.color.rgb = DARK

add_page_number(slide, 17, TOTAL)

# ==================== SLIDE 18: CLOSING ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_body_text(slide, '感谢聆听', Inches(1), Inches(2.0), Inches(11), Inches(1.2), size=52, bold=True, color=WHITE)
add_body_text(slide, 'Tokenhub — 让每一笔 AI 调用都安全、可控、可追溯', Inches(1), Inches(3.5), Inches(11), Inches(0.8), size=20, color=ACCENT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.3), Inches(3), Inches(0.04))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()
add_body_text(slide, '[联系人姓名]  |  [邮箱]  |  [电话]', Inches(1), Inches(4.8), Inches(11), Inches(0.5), size=14, color=GRAY)
add_body_text(slide, '产品 Demo：[URL]  |  GitHub：[URL]', Inches(1), Inches(5.3), Inches(11), Inches(0.5), size=14, color=GRAY)
add_page_number(slide, 18, TOTAL)

# save
output_dir = r'C:\Users\chenjiawei\Desktop\Tokenhub\docs\bp'
os.makedirs(output_dir, exist_ok=True)
ppt_path = os.path.join(output_dir, 'Tokenhub_投资BP_演示文稿.pptx')
prs.save(ppt_path)
print(f'PPT saved: {ppt_path}')
